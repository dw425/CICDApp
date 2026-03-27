"""
PPTX Export for Pipeline Compass.

Generates a 10-slide executive deck from assessment results
using python-pptx. Falls back to text if python-pptx is not available.
"""

import io
from datetime import datetime

from compass.scoring_engine import TIER_COLORS, TIER_LABELS
from compass.antipattern_engine import SEVERITY_COLORS


def generate_pptx_report(assessment: dict, org: dict) -> bytes:
    """
    Generate a PPTX executive deck for an assessment.

    Returns PPTX as bytes for download.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        return _generate_pptx(assessment, org)
    except ImportError:
        # Fallback: return text report as bytes
        return _generate_text_fallback(assessment, org)


def _generate_pptx(assessment: dict, org: dict) -> bytes:
    """Generate PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    composite = assessment.get("composite", {})
    dim_scores = assessment.get("scores", {})
    anti_patterns = assessment.get("anti_patterns", [])
    roadmap = assessment.get("roadmap", {})
    breakdown = composite.get("dimension_breakdown", {})

    overall_score = composite.get("overall_score", 0)
    overall_level = composite.get("overall_level", 1)
    overall_label = composite.get("overall_label", "Initial")
    org_name = org.get("name", "Unknown")
    date_str = datetime.now().strftime("%B %d, %Y")

    bg_color = RGBColor(0x0D, 0x11, 0x17)
    text_color = RGBColor(0xE6, 0xED, 0xF3)
    muted_color = RGBColor(0x8B, 0x94, 0x9E)
    accent_color = RGBColor(0x4B, 0x7B, 0xF5)

    def _hex_to_rgb(hex_str):
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def _add_text(slide, left, top, width, height, text, size=18, color=None, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color or text_color
        p.font.bold = bold
        p.alignment = alignment
        return txBox

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_bg(slide)
    _add_text(slide, 1.5, 1.5, 10, 1, "Pipeline Compass", 40, accent_color, True, PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 2.5, 10, 0.8, "CI/CD Maturity Assessment", 24, text_color, False, PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 4.0, 10, 0.6, org_name, 28, text_color, True, PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 4.8, 10, 0.5, date_str, 16, muted_color, False, PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 6.2, 10, 0.4, "Blueprint Technologies", 14, muted_color, False, PP_ALIGN.CENTER)

    # ── Slide 2: Overall Score ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Overall Maturity Score", 28, accent_color, True)
    oc = _hex_to_rgb(TIER_COLORS.get(overall_level, "#4B7BF5"))
    _add_text(slide, 3.5, 2.0, 6, 1.5, f"{overall_score:.0f}/100", 60, oc, True, PP_ALIGN.CENTER)
    _add_text(slide, 3.5, 3.5, 6, 0.8, f"Level {overall_level}: {overall_label}", 28, text_color, False, PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 5.0, 10, 0.5, f"Weight Profile: {composite.get('weight_profile', 'balanced').replace('_', ' ').title()}", 16, muted_color, False, PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 5.6, 10, 0.5, f"Anti-patterns detected: {len(anti_patterns)}", 16, muted_color, False, PP_ALIGN.CENTER)
    confidence = assessment.get("confidence", {})
    if confidence:
        high_c = sum(1 for v in confidence.values() if v == "high")
        total_c = len(confidence)
        _add_text(slide, 1.5, 6.2, 10, 0.4, f"Confidence: {high_c}/{total_c} dimensions at high confidence", 14, muted_color, False, PP_ALIGN.CENTER)

    # ── Slide 3: Dimension Scorecard ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Dimension Scorecard", 28, accent_color, True)
    y_pos = 1.4
    for dim_id, data in sorted(breakdown.items()):
        name = data.get("display_name", dim_id.replace("_", " ").title())
        score = data.get("score", 0)
        level = data.get("level", 1)
        label = data.get("label", "Initial")
        color = _hex_to_rgb(TIER_COLORS.get(level, "#888"))
        _add_text(slide, 1.0, y_pos, 4, 0.4, name, 14, text_color, True)
        _add_text(slide, 5.5, y_pos, 1.5, 0.4, f"L{level} {label}", 14, color, True)
        _add_text(slide, 7.5, y_pos, 2, 0.4, f"{score:.0f}/100", 14, text_color)
        y_pos += 0.5

    # ── Slide 4: Top 3 Strengths & Gaps ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Strengths & Gaps", 28, accent_color, True)

    sorted_dims = sorted(breakdown.items(), key=lambda x: x[1].get("score", 0), reverse=True)
    _add_text(slide, 1.0, 1.4, 5, 0.5, "Top Strengths", 20, RGBColor(0x34, 0xD3, 0x99), True)
    for i, (dim_id, data) in enumerate(sorted_dims[:3]):
        name = data.get("display_name", dim_id.replace("_", " ").title())
        _add_text(slide, 1.2, 2.0 + i * 0.5, 5, 0.4, f"+ {name} — {data.get('score', 0):.0f}/100", 14, text_color)

    _add_text(slide, 7.0, 1.4, 5, 0.5, "Biggest Gaps", 20, RGBColor(0xEF, 0x44, 0x44), True)
    for i, (dim_id, data) in enumerate(sorted_dims[-3:]):
        name = data.get("display_name", dim_id.replace("_", " ").title())
        _add_text(slide, 7.2, 2.0 + i * 0.5, 5, 0.4, f"- {name} — {data.get('score', 0):.0f}/100", 14, text_color)

    # ── Slide 5: Critical Anti-Patterns ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Critical Anti-Patterns", 28, accent_color, True)
    if anti_patterns:
        for i, ap in enumerate(anti_patterns[:5]):
            sev_color = _hex_to_rgb(SEVERITY_COLORS.get(ap["severity"], "#888"))
            _add_text(slide, 1.0, 1.4 + i * 1.0, 1.5, 0.4, f"[{ap['severity'].upper()}]", 12, sev_color, True)
            _add_text(slide, 2.5, 1.4 + i * 1.0, 8, 0.4, ap["name"], 16, text_color, True)
            _add_text(slide, 2.5, 1.8 + i * 1.0, 8, 0.4, ap.get("recommendation", "")[:120], 11, muted_color)
    else:
        _add_text(slide, 3, 3, 7, 0.6, "No anti-patterns detected!", 20, RGBColor(0x34, 0xD3, 0x99), False, PP_ALIGN.CENTER)

    # ── Slide 6: DORA Metrics ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "DORA 2025 Metrics", 28, accent_color, True)
    dora_data = assessment.get("dora_metrics", {})
    if dora_data:
        y_dora = 1.5
        for metric_name, metric_info in dora_data.items():
            if isinstance(metric_info, dict):
                val = metric_info.get("value", "N/A")
                tier = metric_info.get("tier", "Unknown")
                unit = metric_info.get("unit", "")
                tier_colors_dora = {"Elite": RGBColor(0x3B, 0x82, 0xF6), "High": RGBColor(0x22, 0xC5, 0x5E),
                                     "Medium": RGBColor(0xEA, 0xB3, 0x08), "Low": RGBColor(0xEF, 0x44, 0x44)}
                tc = tier_colors_dora.get(tier, muted_color)
                _add_text(slide, 1.0, y_dora, 4, 0.4, metric_name.replace("_", " ").title(), 16, text_color, True)
                _add_text(slide, 5.5, y_dora, 3, 0.4, f"{val} {unit}", 16, text_color)
                _add_text(slide, 9.0, y_dora, 2, 0.4, tier, 16, tc, True)
                y_dora += 0.6
    else:
        _add_text(slide, 3, 3, 7, 0.6, "No DORA data available — connect data sources to enable", 16, muted_color, False, PP_ALIGN.CENTER)

    # ── Slide 7: Hygiene Summary ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Hygiene Check Summary", 28, accent_color, True)
    hygiene_data = assessment.get("hygiene_scores", [])
    if hygiene_data and isinstance(hygiene_data, list):
        platform_stats = {}
        for h in hygiene_data:
            if not isinstance(h, dict):
                continue
            plat = h.get("platform", "unknown")
            if plat not in platform_stats:
                platform_stats[plat] = {"pass": 0, "warn": 0, "fail": 0}
            status = h.get("status", "unknown")
            if status in platform_stats[plat]:
                platform_stats[plat][status] += 1

        _add_text(slide, 1.0, 1.3, 3, 0.4, "Platform", 14, muted_color, True)
        _add_text(slide, 4.5, 1.3, 1.5, 0.4, "Pass", 14, RGBColor(0x22, 0xC5, 0x5E), True)
        _add_text(slide, 6.0, 1.3, 1.5, 0.4, "Warn", 14, RGBColor(0xEA, 0xB3, 0x08), True)
        _add_text(slide, 7.5, 1.3, 1.5, 0.4, "Fail", 14, RGBColor(0xEF, 0x44, 0x44), True)
        y_hyg = 1.9
        for plat, stats in sorted(platform_stats.items()):
            _add_text(slide, 1.0, y_hyg, 3, 0.4, plat.replace("_", " ").title(), 14, text_color)
            _add_text(slide, 4.5, y_hyg, 1.5, 0.4, str(stats["pass"]), 14, RGBColor(0x22, 0xC5, 0x5E), True)
            _add_text(slide, 6.0, y_hyg, 1.5, 0.4, str(stats["warn"]), 14, RGBColor(0xEA, 0xB3, 0x08), True)
            _add_text(slide, 7.5, y_hyg, 1.5, 0.4, str(stats["fail"]), 14, RGBColor(0xEF, 0x44, 0x44), True)
            y_hyg += 0.5
    else:
        _add_text(slide, 3, 3, 7, 0.6, "No hygiene data — connect CI/CD platforms to enable", 16, muted_color, False, PP_ALIGN.CENTER)

    # ── Slide 8: Quick Wins (was Slide 6) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "30-Day Quick Wins", 28, accent_color, True)
    phases = roadmap.get("phases", [])
    quick_wins = phases[0].get("items", []) if phases else []
    for i, item in enumerate(quick_wins[:5]):
        _add_text(slide, 1.0, 1.4 + i * 1.0, 10, 0.4, item.get("title", ""), 16, text_color, True)
        _add_text(slide, 1.0, 1.8 + i * 1.0, 10, 0.4, f"Effort: {item.get('effort_days', '?')} days | Impact: +{item.get('expected_score_improvement', 0)} points", 12, muted_color)

    # ── Slide 7: Strategic Investments ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "90-Day Strategic Investments", 28, accent_color, True)
    strategic = phases[1].get("items", []) if len(phases) > 1 else []
    for i, item in enumerate(strategic[:5]):
        _add_text(slide, 1.0, 1.4 + i * 1.0, 10, 0.4, item.get("title", ""), 16, text_color, True)
        _add_text(slide, 1.0, 1.8 + i * 1.0, 10, 0.4, f"Effort: {item.get('effort_days', '?')} days | Impact: +{item.get('expected_score_improvement', 0)} points", 12, muted_color)

    # ── Slide 8: ROI Projection ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "ROI Projection", 28, accent_color, True)
    roi = roadmap.get("total_roi_estimate", {})
    roi_cats = roi.get("roi_by_category", {})
    _add_text(slide, 1.5, 1.5, 5, 0.5, "Developer Hours Saved", 16, muted_color)
    _add_text(slide, 1.5, 2.0, 5, 0.6, f"{roi_cats.get('speed', {}).get('hours_saved_annually', 0):,}/year", 24, accent_color, True)
    _add_text(slide, 7.0, 1.5, 5, 0.5, "Incident Reduction", 16, muted_color)
    _add_text(slide, 7.0, 2.0, 5, 0.6, f"{roi_cats.get('quality', {}).get('incident_reduction_pct', 0)}%", 24, RGBColor(0x34, 0xD3, 0x99), True)
    _add_text(slide, 1.5, 3.5, 5, 0.5, "Risk Reduction", 16, muted_color)
    _add_text(slide, 1.5, 4.0, 5, 0.6, f"{roi_cats.get('risk', {}).get('risk_reduction_pct', 0)}%", 24, RGBColor(0xFB, 0xBF, 0x24), True)
    _add_text(slide, 7.0, 3.5, 5, 0.5, "Total Effort", 16, muted_color)
    _add_text(slide, 7.0, 4.0, 5, 0.6, f"{roi.get('total_effort_days', 0)} days", 24, text_color, True)

    # ── Slide 9: Methodology ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Methodology", 28, accent_color, True)
    _add_text(slide, 1.0, 1.5, 10, 0.5, "COMPASS Framework", 20, text_color, True)
    _add_text(slide, 1.0, 2.2, 10, 3, (
        "Pipeline Compass evaluates CI/CD maturity across 9 dimensions using a structured question bank. "
        "Each dimension receives a 0-100 score mapped to L1 (Initial) through L5 (Elite) maturity tiers.\n\n"
        "The composite score uses a weighted geometric mean to prevent high scores in one dimension from "
        "masking critical gaps in another.\n\n"
        "Anti-patterns are detected from response indicators and matched against a library of 12+ known CI/CD "
        "anti-patterns with specific remediation recommendations.\n\n"
        "The improvement roadmap classifies recommendations on an Impact x Effort matrix and assigns them "
        "to 30-day, 90-day, 6-month, and 12-month phases."
    ), 13, muted_color)

    # ── Slide 10: Next Steps ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_text(slide, 0.8, 0.5, 11, 0.6, "Next Steps", 28, accent_color, True)
    _add_text(slide, 1.5, 1.8, 10, 0.5, "1. Review the detailed assessment report", 18, text_color)
    _add_text(slide, 1.5, 2.5, 10, 0.5, "2. Prioritize 30-day quick wins for immediate value", 18, text_color)
    _add_text(slide, 1.5, 3.2, 10, 0.5, "3. Plan 90-day strategic investments with your team", 18, text_color)
    _add_text(slide, 1.5, 3.9, 10, 0.5, "4. Schedule quarterly pulse assessments to track progress", 18, text_color)
    _add_text(slide, 1.5, 4.6, 10, 0.5, "5. Engage Blueprint Technologies for hands-on implementation support", 18, accent_color, True)
    _add_text(slide, 3.5, 6.0, 6, 0.5, "Blueprint Technologies", 14, muted_color, False, PP_ALIGN.CENTER)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _generate_text_fallback(assessment: dict, org: dict) -> bytes:
    """Fallback text-based report when python-pptx is not available."""
    composite = assessment.get("composite", {})
    breakdown = composite.get("dimension_breakdown", {})
    org_name = org.get("name", "Unknown")

    lines = [
        "PIPELINE COMPASS — Executive Deck",
        f"Organization: {org_name}",
        f"Score: {composite.get('overall_score', 0):.0f}/100 — L{composite.get('overall_level', 1)} {composite.get('overall_label', 'Initial')}",
        "",
        "Note: Install python-pptx for proper PPTX generation.",
        "  pip install python-pptx",
        "",
        "Dimension Scores:",
    ]
    for dim_id, data in sorted(breakdown.items()):
        lines.append(f"  {data.get('display_name', dim_id)}: {data.get('score', 0):.0f}/100 (L{data.get('level', 1)})")

    return "\n".join(lines).encode("utf-8")
